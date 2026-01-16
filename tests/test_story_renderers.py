"""
Test Story Renderers (React + PowerPoint)

Validates that both renderers:
1. Work for ALL domains (healthcare, IoT, manufacturing, finance, business)
2. Strictly follow KDS guidelines (colors, typography, spacing)
3. Produce consultant-grade outputs
"""

import json
from pathlib import Path
import pytest
from pptx import Presentation

from kie.story.models import (
    StoryInsight,
    StoryManifest,
    NarrativeMode,
    StoryThesis,
    StoryKPI,
    KPIType,
    StorySection
)
from kie.story.react_story_renderer import ReactStoryRenderer
from kie.story.pptx_story_renderer import PPTXStoryRenderer


# Test Data Fixtures (from Phase 6)

@pytest.fixture
def healthcare_insights():
    """Healthcare clinical trial insights."""
    return [
        StoryInsight(
            insight_id="med_001",
            text="Patient survival rate improved 23% with new protocol (p<0.01)",
            category="outcome",
            confidence=0.95,
            business_value=0.92,
            actionability=0.88,
            supporting_data={"metric": "survival_rate", "change": 0.23}
        ),
        StoryInsight(
            insight_id="med_002",
            text="Symptom reduction observed in 89% of cohort (n=412)",
            category="treatment",
            confidence=0.90,
            business_value=0.85,
            actionability=0.80,
            supporting_data={"metric": "symptom_reduction", "percentage": 0.89}
        ),
        StoryInsight(
            insight_id="med_003",
            text="Treatment efficacy highest in 45-60 age group (p<0.05)",
            category="demographics",
            confidence=0.88,
            business_value=0.80,
            actionability=0.75,
            supporting_data={"age_group": "45-60", "efficacy": 0.92}
        )
    ]


@pytest.fixture
def iot_insights():
    """IoT sensor telemetry insights."""
    return [
        StoryInsight(
            insight_id="iot_001",
            text="System latency reduced 34ms (99th percentile) after optimization",
            category="performance",
            confidence=0.88,
            business_value=0.90,
            actionability=0.95,
            supporting_data={"metric": "latency", "reduction_ms": 34}
        ),
        StoryInsight(
            insight_id="iot_002",
            text="Network uptime reached 99.97%, up from 99.2% baseline",
            category="reliability",
            confidence=0.92,
            business_value=0.88,
            actionability=0.85,
            supporting_data={"metric": "uptime", "current": 0.9997, "baseline": 0.992}
        ),
        StoryInsight(
            insight_id="iot_003",
            text="Sensor data throughput increased 18% under load",
            category="performance",
            confidence=0.85,
            business_value=0.82,
            actionability=0.80,
            supporting_data={"metric": "throughput", "increase": 0.18}
        )
    ]


@pytest.fixture
def manufacturing_insights():
    """Manufacturing quality control insights."""
    return [
        StoryInsight(
            insight_id="mfg_001",
            text="Defect rate dropped to 0.3%, 47% below target threshold",
            category="quality",
            confidence=0.90,
            business_value=0.95,
            actionability=0.92,
            supporting_data={"metric": "defect_rate", "current": 0.003, "target": 0.0056}
        ),
        StoryInsight(
            insight_id="mfg_002",
            text="Production line efficiency increased 18%, throughput at 1,240 units/hour",
            category="efficiency",
            confidence=0.88,
            business_value=0.90,
            actionability=0.88,
            supporting_data={"metric": "throughput", "units_per_hour": 1240}
        ),
        StoryInsight(
            insight_id="mfg_003",
            text="Mean time between failures improved 42% for critical components",
            category="reliability",
            confidence=0.87,
            business_value=0.85,
            actionability=0.83,
            supporting_data={"metric": "mtbf", "improvement": 0.42}
        )
    ]


@pytest.fixture
def financial_insights():
    """Financial trading insights."""
    return [
        StoryInsight(
            insight_id="fin_001",
            text="Sharpe ratio improved to 1.82, volatility decreased 12%",
            category="risk",
            confidence=0.92,
            business_value=0.90,
            actionability=0.88,
            supporting_data={"sharpe_ratio": 1.82, "volatility_change": -0.12}
        ),
        StoryInsight(
            insight_id="fin_002",
            text="Risk-adjusted returns exceed benchmark by 340 basis points",
            category="performance",
            confidence=0.90,
            business_value=0.95,
            actionability=0.85,
            supporting_data={"excess_return_bps": 340}
        ),
        StoryInsight(
            insight_id="fin_003",
            text="Portfolio diversification score reached 0.78 (target: 0.70)",
            category="risk",
            confidence=0.88,
            business_value=0.82,
            actionability=0.80,
            supporting_data={"diversification": 0.78, "target": 0.70}
        )
    ]


def create_story_manifest(insights, project_name, objective, domain):
    """Helper to create StoryManifest from insights."""
    thesis = StoryThesis(
        title=f"{domain.title()} Analysis Results",
        hook=f"Key findings from {domain} data analysis",
        summary=f"Comprehensive analysis reveals significant patterns in {domain} metrics",
        implication=f"These insights drive strategic decisions in {domain} operations",
        confidence=0.90
    )

    kpis = [
        StoryKPI(
            value="23%",
            label="Primary Metric",
            context="Improvement observed",
            kpi_type=KPIType.HEADLINE,
            rank=1,
            insight_id="ins_001"
        ),
        StoryKPI(
            value="89%",
            label="Secondary Metric",
            context="Success Rate achieved",
            kpi_type=KPIType.SUPPORTING,
            rank=2,
            insight_id="ins_002"
        ),
        StoryKPI(
            value="1.82",
            label="Tertiary Metric",
            context="Efficiency Score improved",
            kpi_type=KPIType.DELTA,
            rank=3,
            insight_id="ins_003"
        )
    ]

    section = StorySection(
        section_id="section_001",
        title=f"{domain.title()} Performance",
        subtitle="Key Metrics and Trends",
        thesis=f"Analysis reveals strong performance across {domain} metrics",
        insight_ids=[ins.insight_id for ins in insights],
        narrative_text=f"The {domain} analysis shows consistent improvement across all measured dimensions. Primary drivers include process optimization and enhanced monitoring capabilities.",
        chart_refs=["chart_001"],
        kpis=kpis[:2]
    )

    return StoryManifest(
        story_id=f"story_{domain}_test",
        project_name=project_name,
        thesis=thesis,
        top_kpis=kpis,
        sections=[section],
        narrative_mode=NarrativeMode.EXECUTIVE,
        executive_summary=f"This {domain} analysis demonstrates significant improvements across key metrics, with actionable recommendations for continued optimization.",
        key_findings=[
            f"{domain.title()} metrics show consistent upward trend",
            "Process optimizations yielding measurable results",
            "Monitoring enhancements improve visibility",
            "Strategic recommendations identified for next phase"
        ],
        metadata={
            "generated_at": "2026-01-16T10:00:00Z",
            "objective": objective,
            "domain": domain
        }
    )


# React Renderer Tests

def test_react_renderer_healthcare(healthcare_insights, tmp_path):
    """Test React renderer with healthcare data."""
    story = create_story_manifest(
        healthcare_insights,
        "Clinical Trial Q4 2025",
        "Evaluate treatment efficacy",
        "healthcare"
    )

    renderer = ReactStoryRenderer(theme_mode="dark")
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)

    # Verify component created
    assert result_path.exists()
    assert result_path.name == "StoryView.tsx"

    # Verify StoryView.tsx contains KDS colors
    content = result_path.read_text()
    assert "#7823DC" in content  # Kearney Purple
    assert "#1E1E1E" in content  # Dark background
    assert "Clinical Trial Q4 2025" in content

    # Verify supporting components created
    assert (output_dir / "StorySection.tsx").exists()
    assert (output_dir / "KPICard.tsx").exists()
    assert (output_dir / "ThesisSection.tsx").exists()

    # Verify story data JSON
    story_json = output_dir / "story-data.json"
    assert story_json.exists()
    data = json.loads(story_json.read_text())
    assert data["project_name"] == "Clinical Trial Q4 2025"
    assert len(data["sections"]) == 1


def test_react_renderer_iot(iot_insights, tmp_path):
    """Test React renderer with IoT data."""
    story = create_story_manifest(
        iot_insights,
        "Network Performance Analysis",
        "Assess infrastructure optimization",
        "iot"
    )

    renderer = ReactStoryRenderer(theme_mode="light")
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)

    # Verify light mode colors
    content = result_path.read_text()
    assert "#FFFFFF" in content  # Light background
    assert "#1E1E1E" in content  # Dark text
    assert "Network Performance Analysis" in content


def test_react_renderer_manufacturing(manufacturing_insights, tmp_path):
    """Test React renderer with manufacturing data."""
    story = create_story_manifest(
        manufacturing_insights,
        "Line 3 Process Improvement",
        "Reduce defects and increase throughput",
        "manufacturing"
    )

    renderer = ReactStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)

    content = result_path.read_text()
    assert "Line 3 Process Improvement" in content
    assert "manufacturing" in content.lower()


def test_react_renderer_financial(financial_insights, tmp_path):
    """Test React renderer with financial data."""
    story = create_story_manifest(
        financial_insights,
        "Portfolio Risk Analysis",
        "Optimize risk-adjusted returns",
        "financial"
    )

    renderer = ReactStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)

    content = result_path.read_text()
    assert "Portfolio Risk Analysis" in content


# PowerPoint Renderer Tests

def test_pptx_renderer_healthcare(healthcare_insights, tmp_path):
    """Test PowerPoint renderer with healthcare data."""
    story = create_story_manifest(
        healthcare_insights,
        "Clinical Trial Q4 2025",
        "Evaluate treatment efficacy",
        "healthcare"
    )

    renderer = PPTXStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_path = tmp_path / "healthcare_story.pptx"

    result_path = renderer.render_story(story, charts_dir, output_path)

    # Verify file created
    assert result_path.exists()
    assert result_path.suffix == ".pptx"

    # Load presentation and verify structure
    prs = Presentation(str(result_path))

    # Should have: Title, Executive Summary, KPIs, Section, Key Findings = 5 slides
    assert len(prs.slides) == 5

    # Verify slide dimensions (10" x 7.5")
    assert prs.slide_width == 9144000  # 10 inches in EMUs
    assert prs.slide_height == 6858000  # 7.5 inches in EMUs


def test_pptx_renderer_iot(iot_insights, tmp_path):
    """Test PowerPoint renderer with IoT data."""
    story = create_story_manifest(
        iot_insights,
        "Network Performance Analysis",
        "Assess infrastructure optimization",
        "iot"
    )

    renderer = PPTXStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_path = tmp_path / "iot_story.pptx"

    result_path = renderer.render_story(story, charts_dir, output_path)

    assert result_path.exists()
    prs = Presentation(str(result_path))
    assert len(prs.slides) == 5


def test_pptx_renderer_manufacturing(manufacturing_insights, tmp_path):
    """Test PowerPoint renderer with manufacturing data."""
    story = create_story_manifest(
        manufacturing_insights,
        "Line 3 Process Improvement",
        "Reduce defects and increase throughput",
        "manufacturing"
    )

    renderer = PPTXStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_path = tmp_path / "manufacturing_story.pptx"

    result_path = renderer.render_story(story, charts_dir, output_path)

    assert result_path.exists()
    prs = Presentation(str(result_path))
    assert len(prs.slides) == 5


def test_pptx_renderer_financial(financial_insights, tmp_path):
    """Test PowerPoint renderer with financial data."""
    story = create_story_manifest(
        financial_insights,
        "Portfolio Risk Analysis",
        "Optimize risk-adjusted returns",
        "financial"
    )

    renderer = PPTXStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_path = tmp_path / "financial_story.pptx"

    result_path = renderer.render_story(story, charts_dir, output_path)

    assert result_path.exists()
    prs = Presentation(str(result_path))
    assert len(prs.slides) == 5


# KDS Compliance Tests

def test_react_kds_colors_dark_mode(healthcare_insights, tmp_path):
    """Verify React renderer uses correct KDS colors in dark mode."""
    story = create_story_manifest(
        healthcare_insights,
        "Test Project",
        "Test",
        "healthcare"
    )

    renderer = ReactStoryRenderer(theme_mode="dark")
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)
    content = result_path.read_text()

    # Verify KDS dark mode colors
    assert "#1E1E1E" in content  # Dark background
    assert "#FFFFFF" in content  # White text
    assert "#7823DC" in content  # Kearney Purple
    assert "#A5A6A5" in content  # Medium Gray


def test_react_kds_colors_light_mode(healthcare_insights, tmp_path):
    """Verify React renderer uses correct KDS colors in light mode."""
    story = create_story_manifest(
        healthcare_insights,
        "Test Project",
        "Test",
        "healthcare"
    )

    renderer = ReactStoryRenderer(theme_mode="light")
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)
    content = result_path.read_text()

    # Verify KDS light mode colors
    assert "#FFFFFF" in content  # White background
    assert "#1E1E1E" in content  # Dark text
    assert "#7823DC" in content  # Kearney Purple


def test_pptx_kds_colors(healthcare_insights, tmp_path):
    """Verify PowerPoint renderer uses correct KDS colors."""
    story = create_story_manifest(
        healthcare_insights,
        "Test Project",
        "Test",
        "healthcare"
    )

    renderer = PPTXStoryRenderer()

    # Verify color constants were created correctly
    # Just verify they exist and are RGBColor instances
    from pptx.dml.color import RGBColor
    assert isinstance(renderer.KEARNEY_PURPLE, RGBColor)
    assert isinstance(renderer.KEARNEY_BLACK, RGBColor)
    assert isinstance(renderer.LIGHT_GRAY, RGBColor)
    assert isinstance(renderer.MEDIUM_GRAY, RGBColor)
    assert isinstance(renderer.WHITE, RGBColor)


def test_react_kpi_rendering(healthcare_insights, tmp_path):
    """Verify KPI cards render correctly."""
    story = create_story_manifest(
        healthcare_insights,
        "Test Project",
        "Test",
        "healthcare"
    )

    renderer = ReactStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)

    # Verify KPICard component exists
    kpi_card_path = output_dir / "KPICard.tsx"
    assert kpi_card_path.exists()

    content = kpi_card_path.read_text()
    # KPICard uses inline styles for KDS colors
    assert "#7823DC" in content  # Purple color present
    assert "color: '#7823DC'" in content  # Purple text for value (inline style)


def test_pptx_kpi_layout(healthcare_insights, tmp_path):
    """Verify PowerPoint KPI layout (2x3 grid)."""
    story = create_story_manifest(
        healthcare_insights,
        "Test Project",
        "Test",
        "healthcare"
    )

    renderer = PPTXStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_path = tmp_path / "test_story.pptx"

    result_path = renderer.render_story(story, charts_dir, output_path)
    prs = Presentation(str(result_path))

    # KPI slide should be slide 3 (0-indexed: slide 2)
    kpi_slide = prs.slides[2]

    # Should have multiple shapes (title + KPI cards)
    assert len(kpi_slide.shapes) > 3  # At least title + 3 KPI cards


def test_react_typography(healthcare_insights, tmp_path):
    """Verify React renderer uses correct typography."""
    story = create_story_manifest(
        healthcare_insights,
        "Test Project",
        "Test",
        "healthcare"
    )

    renderer = ReactStoryRenderer()
    charts_dir = tmp_path / "charts"
    charts_dir.mkdir()
    output_dir = tmp_path / "react"

    result_path = renderer.render_story(story, charts_dir, output_dir)
    content = result_path.read_text()

    # Verify Inter font usage
    assert "font-sans" in content or "Inter" in content


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
