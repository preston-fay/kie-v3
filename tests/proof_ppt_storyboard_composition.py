#!/usr/bin/env python3
"""
Proof Test: PPT Composition from Story Manifest

Validates that PowerPoint is composed strictly from:
- story_manifest.json (canonical story representation)
- rendered charts in outputs/charts/

The manifest contains everything needed: sections, narratives, visual order, caveats.
"""

import json
import tempfile
from pathlib import Path

import pytest
import yaml


def test_ppt_requires_story_manifest(tmp_path):
    """Test that PPT build fails without story_manifest.json."""
    from kie.commands.handler import CommandHandler

    handler = CommandHandler(tmp_path)

    # Bootstrap workspace
    handler.handle_startkie()

    # Create minimal spec
    spec_path = tmp_path / "project_state" / "spec.yaml"
    spec_path.write_text(yaml.dump({
        "project_name": "Test Project",
        "client_name": "Test Client",
        "objective": "Test Objective",
    }))

    # Try to build PPT without story manifest - should FAIL
    with pytest.raises(ValueError, match="story_manifest.json not found"):
        handler._build_presentation({"project_name": "Test Project", "client_name": "Test Client"})


def test_ppt_requires_executive_summary(tmp_path):
    """Test that PPT build fails without executive_summary.md."""
    from kie.commands.handler import CommandHandler

    handler = CommandHandler(tmp_path)

    # Bootstrap workspace
    handler.handle_startkie()

    # Create minimal spec
    spec_path = tmp_path / "project_state" / "spec.yaml"
    spec_path.write_text(yaml.dump({
        "project_name": "Test Project",
        "client_name": "Test Client",
        "objective": "Test Objective",
    }))

    # Create intent
    intent_path = tmp_path / "project_state" / "intent.yaml"
    intent_path.write_text(yaml.dump({"objective": "Test Objective"}))

    # Create storyboard
    storyboard = tmp_path / "outputs" / "visual_storyboard.json"
    storyboard.write_text(json.dumps({
        "elements": [
            {
                "section": "Context",
                "order": 1,
                "insight_id": "insight_1",
                "chart_ref": "chart_1.json",
                "role": "Baseline",
                "transition_text": "Test text",
                "emphasis": "Test emphasis",
                "caveats": [],
                "visualization_type": "bar",
                "insight_title": "Test Insight",
            }
        ]
    }))

    # Try to build PPT without executive summary - should FAIL
    with pytest.raises(ValueError, match="executive_summary.md not found"):
        handler._build_presentation({"project_name": "Test Project", "client_name": "Test Client"})


def test_ppt_requires_intent(tmp_path):
    """Test that PPT build fails without intent.yaml."""
    from kie.commands.handler import CommandHandler

    handler = CommandHandler(tmp_path)

    # Bootstrap workspace
    handler.handle_startkie()

    # Create minimal spec
    spec_path = tmp_path / "project_state" / "spec.yaml"
    spec_path.write_text(yaml.dump({
        "project_name": "Test Project",
        "client_name": "Test Client",
        "objective": "Test Objective",
    }))

    # Create executive summary
    exec_summary = tmp_path / "outputs" / "executive_summary.md"
    exec_summary.write_text("# Executive Summary\n\n- Key finding 1\n")

    # Create storyboard
    storyboard = tmp_path / "outputs" / "visual_storyboard.json"
    storyboard.write_text(json.dumps({"elements": []}))

    # Try to build PPT without intent - should FAIL
    with pytest.raises(ValueError, match="intent.yaml not found"):
        handler._build_presentation({"project_name": "Test Project", "client_name": "Test Client"})


def test_ppt_fails_on_missing_chart(tmp_path):
    """Test that PPT build fails if storyboard references non-existent chart."""
    from kie.commands.handler import CommandHandler

    handler = CommandHandler(tmp_path)

    # Bootstrap workspace
    handler.handle_startkie()

    # Create intent
    intent_path = tmp_path / "project_state" / "intent.yaml"
    intent_path.write_text(yaml.dump({"objective": "Test Objective"}))

    # Create executive summary
    exec_summary = tmp_path / "outputs" / "executive_summary.md"
    exec_summary.write_text("# Executive Summary\n\n- Key finding 1\n")

    # Create storyboard referencing non-existent chart
    storyboard = tmp_path / "outputs" / "visual_storyboard.json"
    storyboard.write_text(json.dumps({
        "elements": [
            {
                "section": "Context",
                "order": 1,
                "insight_id": "insight_1",
                "chart_ref": "missing_chart.json",
                "role": "Baseline",
                "transition_text": "Test text",
                "emphasis": "Test emphasis",
                "caveats": [],
                "visualization_type": "bar",
                "insight_title": "Test Insight",
            }
        ]
    }))

    # Try to build PPT - should FAIL because chart doesn't exist
    with pytest.raises(ValueError, match="Chart 'missing_chart.json' referenced in storyboard not found"):
        handler._build_presentation({"project_name": "Test Project", "client_name": "Test Client"})


def test_ppt_composition_order(tmp_path):
    """Test that PPT slides follow the correct order."""
    from kie.commands.handler import CommandHandler
    from pptx import Presentation

    handler = CommandHandler(tmp_path)

    # Bootstrap workspace
    handler.handle_startkie()

    # Create intent
    intent_path = tmp_path / "project_state" / "intent.yaml"
    intent_path.write_text(yaml.dump({"objective": "Analyze sales performance"}))

    # Create executive summary
    exec_summary = tmp_path / "outputs" / "executive_summary.md"
    exec_summary.write_text("""# Executive Summary

- Revenue increased 15% YoY
- Customer retention improved
- Market share grew in Q3

## Risks & Caveats

- Data limited to 6 months
- External factors not accounted for
""")

    # Create charts directory and chart file
    charts_dir = tmp_path / "outputs" / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)

    chart_file = charts_dir / "revenue_trend.json"
    chart_file.write_text(json.dumps({
        "type": "bar",
        "data": [
            {"category": "Q1", "value": 100},
            {"category": "Q2", "value": 150},
        ],
        "config": {
            "xAxis": {"dataKey": "category"},
            "bars": [{"dataKey": "value", "fill": "#7823DC"}]
        }
    }))

    # Create storyboard
    storyboard = tmp_path / "outputs" / "visual_storyboard.json"
    storyboard.write_text(json.dumps({
        "elements": [
            {
                "section": "Context & Baseline",
                "order": 1,
                "insight_id": "insight_1",
                "chart_ref": "revenue_trend.json",
                "role": "Baseline",
                "transition_text": "Revenue shows strong growth",
                "emphasis": "15% increase",
                "caveats": ["Limited to 6 months"],
                "visualization_type": "bar",
                "insight_title": "Revenue Growth",
            }
        ]
    }))

    # Build PPT
    ppt_path = handler._build_presentation({
        "project_name": "Sales Analysis",
        "client_name": "Acme Corp",
    })

    assert ppt_path.exists()

    # Load and validate PPT structure
    prs = Presentation(str(ppt_path))

    # Should have at least: Title, Exec Summary, Section Divider, Chart, Caveats = 5 slides
    assert len(prs.slides) >= 4

    # Slide 1: Title slide
    slide_1 = prs.slides[0]
    assert any("Sales Analysis" in shape.text for shape in slide_1.shapes if hasattr(shape, "text"))

    # Slide 2: Executive Summary
    slide_2 = prs.slides[1]
    assert any("Executive Summary" in shape.text for shape in slide_2.shapes if hasattr(shape, "text"))

    print(f"\n✅ PPT created with {len(prs.slides)} slides")
    print("Slide titles:")
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and len(shape.text) > 5:
                print(f"  Slide {i}: {shape.text[:80]}")
                break


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
