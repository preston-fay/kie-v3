"""
PowerPoint Exporter for Insight Briefs

Converts insight_brief.json to consultant-grade PowerPoint presentations.

HYBRID APPROACH:
- Embeds PNG charts in PPTX (converted from SVG)
- Copies SVG sources to chart_sources/ for consultant editing
"""

import json
import shutil
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from kie.powerpoint.slide_builder import SlideBuilder


class InsightBriefPPTXExporter:
    """
    Export insight briefs to PowerPoint presentations.

    Reads structured JSON from insight_brief.json and generates
    consultant-grade PPTX with KDS compliance.
    """

    def __init__(self, project_root: Path):
        """
        Initialize exporter.

        Args:
            project_root: Project root directory
        """
        self.project_root = project_root

    def export(self, brief_json_path: Path, output_path: Path) -> Path:
        """
        Export brief JSON to PowerPoint.

        Args:
            brief_json_path: Path to insight_brief.json
            output_path: Path to save PPTX file

        Returns:
            Path to generated PPTX file
        """
        # Load brief JSON
        with open(brief_json_path) as f:
            data = json.load(f)

        # Create presentation
        builder = SlideBuilder(title=data.get("strategic_headline", "Insight Brief"))

        # Slide 1: Title
        self._create_title_slide(builder, data)

        # Slide 2: KPI Dashboard
        self._create_kpi_slide(builder, data)

        # Slide 3: Executive Summary
        self._create_executive_summary_slide(builder, data)

        # Slides 4-N: Key Insights (with charts)
        key_insights = data.get("executive_summary", {}).get("key_insights", [])
        charts_dir = self.project_root / "outputs" / "charts"

        for idx, insight in enumerate(key_insights):
            # Try to find matching chart for this insight
            # Charts are named like: insight_0__bar__top_n.svg, insight_2__line.svg
            chart_path = self._find_chart_for_insight(insight, charts_dir, insight_index=idx)
            self._create_insight_slide(builder, insight, chart_path)

        # Slide N+1: Recommendations
        self._create_recommendations_slide(builder, data)

        # Save presentation
        output_path.parent.mkdir(parents=True, exist_ok=True)
        builder.prs.save(str(output_path))

        # Copy SVG sources for consultant editing
        self._copy_svg_sources(charts_dir, output_path.parent)

        return output_path

    def _create_title_slide(self, builder: SlideBuilder, data: dict):
        """Create title slide with strategic headline."""
        builder.add_title_slide(
            title=data.get("strategic_headline", "Insight Brief"),
            subtitle=data.get("business_question", ""),
            date=f"Generated: {data.get('generated_at', '')[:10]}"
        )

    def _create_kpi_slide(self, builder: SlideBuilder, data: dict):
        """Create KPI dashboard slide with text boxes."""
        kpis = data.get("kpis", [])
        if not kpis:
            return

        # Use blank layout
        slide_layout = builder.prs.slide_layouts[6]
        slide = builder.prs.slides.add_slide(slide_layout)
        builder._set_slide_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "📊 Key Performance Indicators"

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(builder.theme.get_text().lstrip("#"))

        # Strategic context (below title)
        context = data.get("strategic_context", "")
        if context:
            context_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(11.333), Inches(0.8)
            )
            tf = context_box.text_frame
            tf.text = context
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(14)
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                builder.theme.get_text("secondary").lstrip("#")
            )

        # KPI cards (grid layout)
        # Display up to 5 KPIs in 2 rows: 3 on top, 2 on bottom
        kpi_width = Inches(3.5)
        kpi_height = Inches(1.5)
        start_y = Inches(2.8)
        spacing = Inches(0.3)

        for i, kpi in enumerate(kpis[:5]):
            # Calculate position (3 per row)
            row = i // 3
            col = i % 3
            x = Inches(1) + col * (kpi_width + spacing)
            y = start_y + row * (kpi_height + spacing)

            # KPI card background (light gray box)
            card = slide.shapes.add_textbox(x, y, kpi_width, kpi_height)
            tf = card.text_frame
            tf.word_wrap = True

            # Value (large text)
            p = tf.add_paragraph()
            p.text = kpi["value"]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string("7823DC")  # Kearney purple

            # Metric name
            p = tf.add_paragraph()
            p.text = kpi["metric"]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                builder.theme.get_text().lstrip("#")
            )

            # Context (smaller)
            context_text = kpi["context"][:50] + "..." if len(kpi["context"]) > 50 else kpi["context"]
            p = tf.add_paragraph()
            p.text = context_text
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                builder.theme.get_text("secondary").lstrip("#")
            )

    def _create_executive_summary_slide(self, builder: SlideBuilder, data: dict):
        """Create executive summary slide."""
        # Use blank layout
        slide_layout = builder.prs.slide_layouts[6]
        slide = builder.prs.slides.add_slide(slide_layout)
        builder._set_slide_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "Executive Summary"

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(builder.theme.get_text().lstrip("#"))

        # Strategic context paragraph
        context = data.get("strategic_context", "")
        if context:
            context_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(11.333), Inches(1.2)
            )
            tf = context_box.text_frame
            tf.text = context
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(16)
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                builder.theme.get_text().lstrip("#")
            )

        # Key insights bullets
        insights = data.get("executive_summary", {}).get("key_insights", [])
        if insights:
            bullets_box = slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11.333), Inches(3.5)
            )
            tf = bullets_box.text_frame
            tf.word_wrap = True

            # Section header
            p = tf.paragraphs[0]
            p.text = "Key Findings:"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(builder.theme.get_text().lstrip("#"))

            # Top 3 insights as bullets
            for insight in insights[:3]:
                p = tf.add_paragraph()
                p.text = insight["headline"]
                p.level = 0
                p.font.size = Pt(16)
                p.font.name = "Inter"
                p.font.color.rgb = RGBColor.from_string(
                    builder.theme.get_text().lstrip("#")
                )

    def _find_chart_for_insight(self, insight: dict, charts_dir: Path, insight_index: int = 0) -> Path | None:
        """
        Find chart SVG file matching this insight.

        Charts are named with insight index pattern: insight_0__bar__top_n.svg

        Args:
            insight: Insight dictionary with headline
            charts_dir: Directory containing SVG chart files
            insight_index: Index of this insight in key_insights list

        Returns:
            Path to matching SVG file, or None if not found
        """
        if not charts_dir.exists():
            return None

        # Look for charts matching this insight index
        # Pattern: insight_{index}__*.svg
        pattern = f"insight_{insight_index}__*.svg"
        matching_charts = list(charts_dir.glob(pattern))

        if matching_charts:
            # If multiple charts for same insight, prefer certain types
            # Priority: line > bar > other
            for chart_path in matching_charts:
                if "__line" in chart_path.stem:
                    return chart_path

            for chart_path in matching_charts:
                if "__bar" in chart_path.stem:
                    return chart_path

            # Return first match if no preferred type
            return matching_charts[0]

        return None

    def _svg_to_png(self, svg_path: Path) -> Path:
        """
        Convert SVG to high-DPI PNG for PowerPoint embedding.

        Args:
            svg_path: Path to SVG file

        Returns:
            Path to generated PNG file
        """
        import subprocess
        import sys

        png_path = svg_path.with_suffix('.png')

        # Use subprocess to ensure environment variables are set before cairocffi loads
        conversion_script = """
import sys
import os

# macOS-specific: Configure cairocffi to find cairo library
if sys.platform == "darwin":
    os.environ['DYLD_FALLBACK_LIBRARY_PATH'] = '/opt/homebrew/lib'

try:
    import cairosvg
except ImportError:
    print("ERROR: cairosvg is required. Install with: pip install cairosvg", file=sys.stderr)
    sys.exit(1)

import sys
svg_path = sys.argv[1]
png_path = sys.argv[2]

cairosvg.svg2png(
    url=svg_path,
    write_to=png_path,
    output_width=1600,
    dpi=300
)
"""

        result = subprocess.run(
            [sys.executable, "-c", conversion_script, str(svg_path), str(png_path)],
            capture_output=True,
            text=True
        )

        if result.returncode != 0:
            raise RuntimeError(f"SVG conversion failed: {result.stderr}")

        return png_path

    def _copy_svg_sources(self, charts_dir: Path, output_dir: Path):
        """
        Copy SVG source files to chart_sources/ for consultant editing.

        Args:
            charts_dir: Directory containing source SVG files
            output_dir: Output directory (parent of chart_sources/)
        """
        if not charts_dir.exists():
            return

        # Create chart_sources directory
        sources_dir = output_dir / "chart_sources"
        sources_dir.mkdir(parents=True, exist_ok=True)

        # Copy all SVG files
        svg_files = list(charts_dir.glob("*.svg"))
        for svg_file in svg_files:
            dest = sources_dir / svg_file.name
            shutil.copy2(svg_file, dest)

        if svg_files:
            print(f"  → Copied {len(svg_files)} SVG sources to chart_sources/")

    def _create_insight_slide(self, builder: SlideBuilder, insight: dict, chart_path: Path | None = None):
        """Create slide for a single key insight with optional chart."""
        # Use blank layout
        slide_layout = builder.prs.slide_layouts[6]
        slide = builder.prs.slides.add_slide(slide_layout)
        builder._set_slide_background(slide)

        # Determine layout: if chart exists, use two-column layout
        has_chart = chart_path is not None and chart_path.exists()
        text_width = Inches(6) if has_chart else Inches(11.333)

        # Title (insight headline)
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = insight["headline"]
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(builder.theme.get_text().lstrip("#"))

        # Supporting text (narrower if chart present)
        supporting_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), text_width, Inches(1.2)
        )
        tf = supporting_box.text_frame
        tf.text = insight["supporting_text"]
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(16)
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(
            builder.theme.get_text().lstrip("#")
        )

        # Chart embedding (if available)
        if has_chart:
            try:
                # Convert SVG to PNG
                png_path = self._svg_to_png(chart_path)

                # Add chart to right side of slide
                chart_left = Inches(7.5)
                chart_top = Inches(1.5)
                chart_width = Inches(5)
                chart_height = Inches(5.5)

                slide.shapes.add_picture(
                    str(png_path),
                    chart_left,
                    chart_top,
                    width=chart_width,
                    height=chart_height
                )
            except Exception as e:
                # If chart embedding fails, continue without it
                print(f"  ⚠️ Chart embedding failed: {e}")

        # "So What" callout box (adjust width if chart present)
        so_what = insight.get("so_what", "")
        if so_what:
            callout_box = slide.shapes.add_textbox(
                Inches(1), Inches(3), text_width, Inches(1.2)
            )
            tf = callout_box.text_frame
            tf.word_wrap = True

            # Header
            p = tf.paragraphs[0]
            p.text = "💡 So What"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string("7823DC")  # Purple

            # Content
            p = tf.add_paragraph()
            p.text = so_what
            p.font.size = Pt(14)
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                builder.theme.get_text().lstrip("#")
            )

        # Evidence bullets (adjust width if chart present)
        evidence = insight.get("evidence", [])
        if evidence:
            evidence_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), text_width, Inches(2.5)
            )
            tf = evidence_box.text_frame
            tf.word_wrap = True

            # Header
            p = tf.paragraphs[0]
            p.text = "Evidence:"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(builder.theme.get_text().lstrip("#"))

            # Evidence bullets (top 3)
            for ev in evidence[:3]:
                p = tf.add_paragraph()
                label = ev.get("label", ev.get("reference", ""))
                value = ev.get("value", "")
                p.text = f"{label}: {value}"
                p.level = 0
                p.font.size = Pt(14)
                p.font.name = "Inter"
                p.font.color.rgb = RGBColor.from_string(
                    builder.theme.get_text("secondary").lstrip("#")
                )

    def _create_recommendations_slide(self, builder: SlideBuilder, data: dict):
        """Create recommendations slide."""
        recommendations = data.get("recommendations", [])
        if not recommendations:
            return

        # Use blank layout
        slide_layout = builder.prs.slide_layouts[6]
        slide = builder.prs.slides.add_slide(slide_layout)
        builder._set_slide_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "Recommended Next Actions"

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(builder.theme.get_text().lstrip("#"))

        # Recommendations bullets
        bullets_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11.333), Inches(5.5)
        )
        tf = bullets_box.text_frame
        tf.word_wrap = True

        for i, rec in enumerate(recommendations[:5], 1):
            # Extract text (handle both catalog format and generated format)
            if isinstance(rec, dict):
                rec_text = rec.get("text") or f"{rec.get('headline', '')}: {rec.get('supporting_text', '')}"
            else:
                rec_text = str(rec)

            # Clean markdown formatting (remove ** bold **)
            rec_text = rec_text.replace("**", "")

            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = rec_text
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                builder.theme.get_text().lstrip("#")
            )
