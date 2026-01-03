"""
KIE Slide Engine

Brand-compliant PowerPoint presentation generation with design system integration.
Supports dark mode, accessibility, and consulting-style layouts.

Usage:
    from kie.slides import Presentation

    pres = Presentation()
    pres.add_title("Revenue Analysis", "Q4 2024")
    pres.add_content("Key Findings", ["Revenue up 15%", "Costs down 10%"])
    pres.add_chart("Trend Analysis", "outputs/charts/trend.png")
    pres.save("exports/presentation.pptx")
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from kie.brand import DesignSystem, load_design_system
from kie.slides.types import SlideType, SlideLayout

logger = logging.getLogger(__name__)

# Try to import python-pptx
try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    PptxPresentation = None
    RGBColor = None


@dataclass
class SlideConfig:
    """Configuration for slide generation."""
    # Dimensions (16:9 widescreen)
    width: float = 13.333  # inches
    height: float = 7.5    # inches

    # Dark mode (default per KDS)
    dark_mode: bool = True

    # Design system slug
    design_system: str = "kearney"

    # Content limits
    max_bullets_per_slide: int = 6
    max_table_rows: int = 10

    # Typography
    title_font_size: int = 32
    subtitle_font_size: int = 20
    body_font_size: int = 18
    caption_font_size: int = 12

    # Spacing (inches)
    margin: float = 0.75
    title_top: float = 0.5
    content_top: float = 1.5

    # Source attribution
    include_source: bool = True
    source_position: str = "bottom_right"


class Presentation:
    """
    Brand-compliant PowerPoint presentation generator.

    Integrates with KIE design system for consistent branding.
    Supports method chaining for fluent API.
    """

    def __init__(
        self,
        config: Optional[SlideConfig] = None,
        design_system: Optional[Union[str, DesignSystem]] = None,
    ):
        """
        Initialize presentation.

        Args:
            config: Slide configuration. Defaults to standard KDS settings.
            design_system: Design system slug or instance. Defaults to "kearney".

        Raises:
            ImportError: If python-pptx is not installed.
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for slide generation. "
                "Install with: pip install python-pptx"
            )

        self.config = config or SlideConfig()

        # Load design system
        if design_system is None:
            self.ds = load_design_system(self.config.design_system)
        elif isinstance(design_system, str):
            self.ds = load_design_system(design_system)
        else:
            self.ds = design_system

        # Create presentation
        self._prs = PptxPresentation()
        self._prs.slide_width = Inches(self.config.width)
        self._prs.slide_height = Inches(self.config.height)

        self._slide_count = 0
        self._sections: List[str] = []

        # Cache RGB colors
        self._colors = self._build_color_cache()

    def _build_color_cache(self) -> Dict[str, "RGBColor"]:
        """Build RGB color cache from design system."""
        def hex_to_rgb_color(hex_color: str) -> "RGBColor":
            hex_color = hex_color.lstrip("#")
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)

        return {
            "primary": hex_to_rgb_color(self.ds.primary),
            "primary_light": hex_to_rgb_color(self.ds.primary_light),
            "primary_dark": hex_to_rgb_color(self.ds.primary_dark),
            "background": hex_to_rgb_color(
                self.ds.background_dark if self.config.dark_mode else self.ds.background_light
            ),
            "text": hex_to_rgb_color(
                self.ds.text_light if self.config.dark_mode else self.ds.text_dark
            ),
            "text_muted": hex_to_rgb_color("#999999"),
            "alt_background": hex_to_rgb_color("#2A2A2A" if self.config.dark_mode else "#F5F5F5"),
        }

    def _set_slide_background(self, slide) -> None:
        """Set slide background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._colors["background"]

    def _add_textbox(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int = 18,
        bold: bool = False,
        alignment: "PP_ALIGN" = None,
        color: Optional["RGBColor"] = None,
    ):
        """Add a text box to slide."""
        if alignment is None:
            alignment = PP_ALIGN.LEFT

        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self._colors["text"]
        p.font.name = self.ds.font_primary
        p.alignment = alignment

        return textbox

    def _add_accent_line(self, slide, left: float, top: float, width: float) -> None:
        """Add a brand accent line."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.05)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._colors["primary"]
        shape.line.fill.background()

    def add_title(
        self,
        title: str,
        subtitle: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a title slide.

        Args:
            title: Main title text
            subtitle: Optional subtitle

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]  # Blank
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title - centered
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=2.5,
            width=self.config.width - (2 * self.config.margin),
            height=1.5,
            text=title,
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Accent line
        line_width = 4.33
        self._add_accent_line(
            slide,
            left=(self.config.width - line_width) / 2,
            top=3.9,
            width=line_width,
        )

        # Subtitle
        if subtitle:
            self._add_textbox(
                slide,
                left=self.config.margin,
                top=4.0,
                width=self.config.width - (2 * self.config.margin),
                height=0.75,
                text=subtitle,
                font_size=self.config.subtitle_font_size,
                bold=False,
                alignment=PP_ALIGN.CENTER,
            )

        self._slide_count += 1
        logger.debug(f"Added title slide: {title}")
        return self

    def add_section(
        self,
        title: str,
        section_number: Optional[int] = None,
    ) -> "Presentation":
        """
        Add a section divider slide.

        Args:
            title: Section title
            section_number: Optional section number to display

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Section number (large, purple)
        if section_number is not None:
            num_str = f"0{section_number}" if section_number < 10 else str(section_number)
            num_box = self._add_textbox(
                slide,
                left=self.config.margin,
                top=2.5,
                width=self.config.width - (2 * self.config.margin),
                height=1.0,
                text=num_str,
                font_size=72,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                color=self._colors["primary"],
            )

        # Section title
        title_top = 3.5 if section_number else 3.0
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=title_top,
            width=self.config.width - (2 * self.config.margin),
            height=1.0,
            text=title,
            font_size=36,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        self._sections.append(title)
        self._slide_count += 1
        logger.debug(f"Added section slide: {title}")
        return self

    def add_content(
        self,
        title: str,
        bullets: List[str],
        subtitle: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title (action title format recommended)
            bullets: List of bullet point strings
            subtitle: Optional subtitle below title

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=self.config.title_top,
            width=self.config.width - (2 * self.config.margin),
            height=0.75,
            text=title,
            font_size=self.config.title_font_size,
            bold=True,
        )

        # Subtitle
        content_top = self.config.content_top
        if subtitle:
            self._add_textbox(
                slide,
                left=self.config.margin,
                top=1.25,
                width=self.config.width - (2 * self.config.margin),
                height=0.5,
                text=subtitle,
                font_size=self.config.body_font_size,
            )
            content_top = 1.75

        # Bullet points
        bullet_box = slide.shapes.add_textbox(
            Inches(self.config.margin),
            Inches(content_top),
            Inches(self.config.width - (2 * self.config.margin)),
            Inches(5.25),
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullets[:self.config.max_bullets_per_slide]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2022  {point}"
            p.font.size = Pt(self.config.subtitle_font_size)
            p.font.color.rgb = self._colors["text"]
            p.font.name = self.ds.font_primary
            p.space_after = Pt(12)

        self._slide_count += 1
        logger.debug(f"Added content slide: {title}")
        return self

    def add_chart(
        self,
        title: str,
        chart_path: Union[str, Path],
        caption: Optional[str] = None,
        source: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a slide with a chart image.

        Args:
            title: Slide title
            chart_path: Path to chart image file
            caption: Optional caption below chart
            source: Optional data source attribution

        Returns:
            self for method chaining

        Raises:
            FileNotFoundError: If chart file doesn't exist
        """
        chart_path = Path(chart_path)
        if not chart_path.exists():
            raise FileNotFoundError(f"Chart not found: {chart_path}")

        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=self.config.title_top,
            width=self.config.width - (2 * self.config.margin),
            height=0.75,
            text=title,
            font_size=self.config.title_font_size,
            bold=True,
        )

        # Chart image (centered)
        chart_width = 10.0
        chart_height = 5.5
        chart_left = (self.config.width - chart_width) / 2
        chart_top = 1.25

        slide.shapes.add_picture(
            str(chart_path),
            Inches(chart_left),
            Inches(chart_top),
            Inches(chart_width),
            Inches(chart_height),
        )

        # Caption
        if caption:
            self._add_textbox(
                slide,
                left=self.config.margin,
                top=6.75,
                width=self.config.width - (2 * self.config.margin),
                height=0.5,
                text=caption,
                font_size=self.config.caption_font_size,
                alignment=PP_ALIGN.CENTER,
            )

        # Source attribution
        if source and self.config.include_source:
            self._add_textbox(
                slide,
                left=self.config.margin,
                top=6.95,
                width=self.config.width - (2 * self.config.margin),
                height=0.3,
                text=f"Source: {source}",
                font_size=10,
                alignment=PP_ALIGN.RIGHT,
                color=self._colors["text_muted"],
            )

        self._slide_count += 1
        logger.debug(f"Added chart slide: {title}")
        return self

    def add_two_column(
        self,
        title: str,
        left_content: List[str],
        right_content: List[str],
        left_header: Optional[str] = None,
        right_header: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a two-column content slide.

        Args:
            title: Slide title
            left_content: Bullet points for left column
            right_content: Bullet points for right column
            left_header: Optional header for left column
            right_header: Optional header for right column

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=self.config.title_top,
            width=self.config.width - (2 * self.config.margin),
            height=0.75,
            text=title,
            font_size=self.config.title_font_size,
            bold=True,
        )

        col_width = 5.5
        col_left = self.config.margin
        col_right = 7.0

        # Left column
        self._add_column_content(
            slide, col_left, left_header, left_content, col_width
        )

        # Right column
        self._add_column_content(
            slide, col_right, right_header, right_content, col_width
        )

        self._slide_count += 1
        logger.debug(f"Added two-column slide: {title}")
        return self

    def _add_column_content(
        self,
        slide,
        left: float,
        header: Optional[str],
        content: List[str],
        width: float,
    ) -> None:
        """Add column content (header + bullets)."""
        top = 1.5

        if header:
            header_box = self._add_textbox(
                slide,
                left=left,
                top=top,
                width=width,
                height=0.5,
                text=header,
                font_size=self.config.subtitle_font_size,
                bold=True,
                color=self._colors["primary"],
            )
            top = 2.0

        content_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2022  {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self._colors["text"]
            p.font.name = self.ds.font_primary
            p.space_after = Pt(8)

    def add_table(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
    ) -> "Presentation":
        """
        Add a slide with a data table.

        Args:
            title: Slide title
            headers: Column headers
            rows: Table rows (list of lists)

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=self.config.title_top,
            width=self.config.width - (2 * self.config.margin),
            height=0.75,
            text=title,
            font_size=self.config.title_font_size,
            bold=True,
        )

        # Table
        n_rows = min(len(rows), self.config.max_table_rows) + 1
        n_cols = len(headers)

        table_width = 11.5
        table_height = 0.4 * n_rows
        col_width = table_width / n_cols

        table = slide.shapes.add_table(
            rows=n_rows,
            cols=n_cols,
            left=Inches(0.9),
            top=Inches(1.5),
            width=Inches(table_width),
            height=Inches(table_height),
        ).table

        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._colors["primary"]

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = self.ds.font_primary
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, row in enumerate(rows[:self.config.max_table_rows], start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)

                # Alternate row colors
                cell.fill.solid()
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = self._colors["alt_background"]
                else:
                    cell.fill.fore_color.rgb = self._colors["background"]

                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = self._colors["text"]
                p.font.name = self.ds.font_primary
                p.alignment = PP_ALIGN.CENTER

        self._slide_count += 1
        logger.debug(f"Added table slide: {title}")
        return self

    def add_key_takeaway(
        self,
        message: str,
        supporting_text: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a key takeaway slide with prominent message.

        Args:
            message: Main takeaway message (displayed large)
            supporting_text: Optional supporting detail

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Large message
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=2.5,
            width=self.config.width - (2 * self.config.margin),
            height=2.0,
            text=message,
            font_size=36,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            color=self._colors["primary"],
        )

        # Supporting text
        if supporting_text:
            self._add_textbox(
                slide,
                left=1.5,
                top=4.5,
                width=self.config.width - 3.0,
                height=1.0,
                text=supporting_text,
                font_size=self.config.body_font_size,
                alignment=PP_ALIGN.CENTER,
            )

        self._slide_count += 1
        logger.debug(f"Added key takeaway slide")
        return self

    def add_quote(
        self,
        quote: str,
        attribution: str,
        title: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a quote slide.

        Args:
            quote: The quote text
            attribution: Who said it (name, title)
            title: Optional slide title

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Optional title
        if title:
            self._add_textbox(
                slide,
                left=self.config.margin,
                top=self.config.title_top,
                width=self.config.width - (2 * self.config.margin),
                height=0.75,
                text=title,
                font_size=self.config.title_font_size,
                bold=True,
            )

        # Quote with quotation marks
        quote_top = 2.0 if title else 2.5
        self._add_textbox(
            slide,
            left=1.5,
            top=quote_top,
            width=self.config.width - 3.0,
            height=2.5,
            text=f'"{quote}"',
            font_size=28,
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )

        # Attribution
        self._add_textbox(
            slide,
            left=1.5,
            top=5.0,
            width=self.config.width - 3.0,
            height=0.75,
            text=f"- {attribution}",
            font_size=self.config.body_font_size,
            alignment=PP_ALIGN.CENTER,
            color=self._colors["text_muted"],
        )

        self._slide_count += 1
        logger.debug(f"Added quote slide")
        return self

    def add_agenda(
        self,
        title: str,
        items: List[str],
        current_item: Optional[int] = None,
    ) -> "Presentation":
        """
        Add an agenda slide.

        Args:
            title: Slide title (e.g., "Today's Agenda")
            items: List of agenda items
            current_item: Optional index (0-based) to highlight

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=self.config.title_top,
            width=self.config.width - (2 * self.config.margin),
            height=0.75,
            text=title,
            font_size=self.config.title_font_size,
            bold=True,
        )

        # Agenda items with numbers
        item_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.75),
            Inches(self.config.width - 3.0), Inches(4.5)
        )
        tf = item_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i + 1}.  {item}"
            p.font.size = Pt(22)
            p.font.name = self.ds.font_primary
            p.space_after = Pt(16)

            # Highlight current item
            if current_item is not None and i == current_item:
                p.font.bold = True
                p.font.color.rgb = self._colors["primary"]
            else:
                p.font.color.rgb = self._colors["text"]

        self._slide_count += 1
        logger.debug(f"Added agenda slide: {title}")
        return self

    def add_closing(
        self,
        title: str = "Thank You",
        contact_info: Optional[List[str]] = None,
    ) -> "Presentation":
        """
        Add a closing slide.

        Args:
            title: Closing title (default: "Thank You")
            contact_info: Optional contact information lines

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title in brand color
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=2.5,
            width=self.config.width - (2 * self.config.margin),
            height=1.5,
            text=title,
            font_size=44,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            color=self._colors["primary"],
        )

        # Contact info
        if contact_info:
            info_box = slide.shapes.add_textbox(
                Inches(self.config.margin),
                Inches(4.0),
                Inches(self.config.width - (2 * self.config.margin)),
                Inches(2.0),
            )
            tf = info_box.text_frame
            tf.word_wrap = True

            for i, line in enumerate(contact_info):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = self._colors["text"]
                p.font.name = self.ds.font_primary
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(6)

        self._slide_count += 1
        logger.debug(f"Added closing slide: {title}")
        return self

    def add_placeholder(
        self,
        title: str,
        placeholder_text: str,
        instructions: Optional[str] = None,
    ) -> "Presentation":
        """
        Add a placeholder slide for manual content.

        Args:
            title: Slide title
            placeholder_text: Description of what should go here
            instructions: Optional instructions for creating content

        Returns:
            self for method chaining
        """
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide)

        # Title
        self._add_textbox(
            slide,
            left=self.config.margin,
            top=self.config.title_top,
            width=self.config.width - (2 * self.config.margin),
            height=0.75,
            text=title,
            font_size=self.config.title_font_size,
            bold=True,
        )

        # Placeholder box
        box_width = 11.0
        box_height = 5.0
        box_left = (self.config.width - box_width) / 2

        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(box_left), Inches(1.5),
            Inches(box_width), Inches(box_height)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = self._colors["alt_background"]
        placeholder.line.color.rgb = self._colors["primary"]
        placeholder.line.width = Pt(2)

        # Placeholder text
        tf = placeholder.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"[PLACEHOLDER: {placeholder_text}]"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self._colors["primary"]
        p.alignment = PP_ALIGN.CENTER

        if instructions:
            p2 = tf.add_paragraph()
            p2.text = f"\n{instructions}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = self._colors["text_muted"]
            p2.alignment = PP_ALIGN.CENTER

        # Add to notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            f"PLACEHOLDER SLIDE\n\n"
            f"Content needed: {placeholder_text}\n\n"
            f"Instructions: {instructions or 'Replace with final content'}"
        )

        self._slide_count += 1
        logger.debug(f"Added placeholder slide: {title}")
        return self

    @property
    def slide_count(self) -> int:
        """Get the number of slides."""
        return self._slide_count

    @property
    def sections(self) -> List[str]:
        """Get list of section titles."""
        return self._sections.copy()

    def save(self, path: Union[str, Path]) -> Path:
        """
        Save the presentation.

        Args:
            path: Output file path

        Returns:
            Path to saved file
        """
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)

        self._prs.save(str(path))
        logger.info(f"Presentation saved to {path} ({self._slide_count} slides)")
        return path

    def close(self) -> None:
        """Close the presentation without saving."""
        self._prs = None
        logger.debug("Presentation closed without saving")
