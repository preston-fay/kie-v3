"""
Slide Builder with KDS Templates

Creates KDS-compliant PowerPoint slides with standard layouts.
"""

from typing import Optional, List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

from kie.brand.theme import get_theme
from kie.powerpoint.chart_embedder import PowerPointChartEmbedder


class SlideBuilder:
    """
    Build PowerPoint presentations with KDS styling.

    Provides standard slide layouts and templates.
    """

    def __init__(self, title: str = "Presentation"):
        """
        Initialize slide builder.

        Args:
            title: Presentation title
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 widescreen
        self.prs.slide_height = Inches(7.5)

        self.theme = get_theme()
        self.chart_embedder = PowerPointChartEmbedder()
        self.title = title

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        author: Optional[str] = None,
        date: Optional[str] = None,
    ):
        """
        Add title slide.

        Args:
            title: Main title
            subtitle: Subtitle
            author: Author name
            date: Date string

        Returns:
            Slide object
        """
        # Use blank layout
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Set background
        self._set_slide_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(self.theme.get_text().lstrip("#"))

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(11.333), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.text = subtitle

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                self.theme.get_text("secondary").lstrip("#")
            )

        # Author and date (bottom right)
        if author or date:
            info_text = []
            if author:
                info_text.append(author)
            if date:
                info_text.append(date)

            info_box = slide.shapes.add_textbox(
                Inches(9), Inches(6.8), Inches(3.333), Inches(0.5)
            )
            tf = info_box.text_frame
            tf.text = " | ".join(info_text)

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(12)
            p.font.name = "Inter"
            p.font.color.rgb = RGBColor.from_string(
                self.theme.get_text("tertiary").lstrip("#")
            )

        # Kearney branding (bottom left)
        self._add_branding(slide)

        return slide

    def add_section_slide(self, section_title: str):
        """
        Add section divider slide.

        Args:
            section_title: Section title

        Returns:
            Slide object
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide)

        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = section_title

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(self.theme.colors.brand_primary.lstrip("#"))

        # Accent line
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(1),
            Inches(4.5),
            Inches(3),
            Inches(0),
        )
        line.line.color.rgb = RGBColor.from_string(self.theme.colors.brand_primary.lstrip("#"))
        line.line.width = Pt(3)

        self._add_branding(slide)

        return slide

    def add_content_slide(
        self,
        title: str,
        bullet_points: Optional[List[str]] = None,
        notes: Optional[str] = None,
    ):
        """
        Add content slide with bullet points.

        Args:
            title: Slide title
            bullet_points: List of bullet points
            notes: Speaker notes

        Returns:
            Slide object
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide)
        self._add_slide_title(slide, title)

        # Bullet points
        if bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8), Inches(11.333), Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullet_points):
                if i > 0:
                    tf.add_paragraph()

                p = tf.paragraphs[i]
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.name = "Inter"
                p.font.color.rgb = RGBColor.from_string(self.theme.get_text().lstrip("#"))

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        self._add_branding(slide)
        self._add_slide_number(slide)

        return slide

    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: List[Dict[str, Any]],
        notes: Optional[str] = None,
        **chart_kwargs,
    ):
        """
        Add slide with embedded chart.

        Args:
            title: Slide title
            chart_type: Type of chart ('bar', 'line', 'pie', 'area')
            data: Chart data
            notes: Speaker notes
            **chart_kwargs: Chart-specific parameters

        Returns:
            Slide object
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide)
        self._add_slide_title(slide, title)

        # Embed chart
        chart_position = (1, 2, 11.333, 5)  # Leave room for title and branding

        if chart_type == "bar":
            self.chart_embedder.embed_bar_chart(
                slide, data, position=chart_position, **chart_kwargs
            )
        elif chart_type == "line":
            self.chart_embedder.embed_line_chart(
                slide, data, position=chart_position, **chart_kwargs
            )
        elif chart_type == "pie":
            self.chart_embedder.embed_pie_chart(
                slide, data, position=chart_position, **chart_kwargs
            )
        elif chart_type == "area":
            self.chart_embedder.embed_area_chart(
                slide, data, position=chart_position, **chart_kwargs
            )

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        self._add_branding(slide)
        self._add_slide_number(slide)

        return slide

    def add_two_chart_slide(
        self,
        title: str,
        chart1_type: str,
        chart1_data: List[Dict[str, Any]],
        chart2_type: str,
        chart2_data: List[Dict[str, Any]],
        notes: Optional[str] = None,
        **kwargs,
    ):
        """
        Add slide with two side-by-side charts.

        Args:
            title: Slide title
            chart1_type: First chart type
            chart1_data: First chart data
            chart2_type: Second chart type
            chart2_data: Second chart data
            notes: Speaker notes

        Returns:
            Slide object
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._set_slide_background(slide)
        self._add_slide_title(slide, title)

        # Left chart
        chart1_kwargs = {k.replace("chart1_", ""): v for k, v in kwargs.items() if k.startswith("chart1_")}
        chart1_position = (0.5, 2, 5.9, 4.5)

        if chart1_type == "bar":
            self.chart_embedder.embed_bar_chart(slide, chart1_data, position=chart1_position, **chart1_kwargs)
        elif chart1_type == "line":
            self.chart_embedder.embed_line_chart(slide, chart1_data, position=chart1_position, **chart1_kwargs)
        elif chart1_type == "pie":
            self.chart_embedder.embed_pie_chart(slide, chart1_data, position=chart1_position, **chart1_kwargs)

        # Right chart
        chart2_kwargs = {k.replace("chart2_", ""): v for k, v in kwargs.items() if k.startswith("chart2_")}
        chart2_position = (6.9, 2, 5.9, 4.5)

        if chart2_type == "bar":
            self.chart_embedder.embed_bar_chart(slide, chart2_data, position=chart2_position, **chart2_kwargs)
        elif chart2_type == "line":
            self.chart_embedder.embed_line_chart(slide, chart2_data, position=chart2_position, **chart2_kwargs)
        elif chart2_type == "pie":
            self.chart_embedder.embed_pie_chart(slide, chart2_data, position=chart2_position, **chart2_kwargs)

        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        self._add_branding(slide)
        self._add_slide_number(slide)

        return slide

    def save(self, output_path: str) -> Path:
        """
        Save presentation.

        Args:
            output_path: Output file path

        Returns:
            Path to saved file
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        return output_path

    def _set_slide_background(self, slide):
        """Set KDS background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self.theme.get_background().lstrip("#"))

    def _add_slide_title(self, slide, title: str):
        """Add title to slide."""
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.333), Inches(1)
        )
        tf = title_box.text_frame
        tf.text = title

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(self.theme.get_text().lstrip("#"))

    def _add_branding(self, slide):
        """Add Kearney branding (logo placeholder)."""
        # Logo placeholder (bottom left)
        logo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8), Inches(2), Inches(0.5)
        )
        tf = logo_box.text_frame
        tf.text = "KEARNEY"

        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(self.theme.colors.brand_primary.lstrip("#"))

    def _add_slide_number(self, slide):
        """Add slide number (bottom right)."""
        slide_num = len(self.prs.slides)

        num_box = slide.shapes.add_textbox(
            Inches(12.333), Inches(6.8), Inches(0.5), Inches(0.5)
        )
        tf = num_box.text_frame
        tf.text = str(slide_num)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(10)
        p.font.name = "Inter"
        p.font.color.rgb = RGBColor.from_string(self.theme.get_text("tertiary").lstrip("#"))
