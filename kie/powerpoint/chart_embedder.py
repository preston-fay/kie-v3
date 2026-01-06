"""
PowerPoint Native Chart Embedder

Embeds native, editable charts in PowerPoint slides (not images).
Charts remain fully editable after generation.
"""

from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from kie.brand.colors import KDSColors
from kie.brand.theme import get_theme


class PowerPointChartEmbedder:
    """
    Embed native, editable charts in PowerPoint.

    Charts are native PowerPoint objects, not images.
    Users can double-click to edit data and formatting.
    """

    def __init__(self):
        """Initialize chart embedder."""
        self.theme = get_theme()
        self.colors = KDSColors.get_chart_colors(10)

    def embed_bar_chart(
        self,
        slide,
        data: list[dict[str, Any]],
        x_key: str,
        y_keys: list[str],
        title: str | None = None,
        position: tuple[float, float, float, float] = None,
        stacked: bool = False,
    ):
        """
        Embed native bar chart in slide.

        Args:
            slide: PowerPoint slide object
            data: List of data dicts
            x_key: Key for x-axis (categories)
            y_keys: Keys for y-axis (series)
            title: Chart title
            position: (left, top, width, height) in inches
            stacked: Whether to stack bars

        Returns:
            Chart object
        """
        # Default position (centered, large)
        if position is None:
            position = (0.5, 1.5, 12, 5)

        left, top, width, height = [Inches(x) for x in position]

        # Determine chart type
        chart_type = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = [item[x_key] for item in data]

        # Add series
        for _i, y_key in enumerate(y_keys):
            values = [item[y_key] for item in data]
            chart_data.add_series(y_key, values)

        # Add chart to slide
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart

        # Apply KDS styling
        self._style_chart(chart, title, len(y_keys))

        return chart

    def embed_line_chart(
        self,
        slide,
        data: list[dict[str, Any]],
        x_key: str,
        y_keys: list[str],
        title: str | None = None,
        position: tuple[float, float, float, float] = None,
        smooth: bool = True,
    ):
        """
        Embed native line chart in slide.

        Args:
            slide: PowerPoint slide object
            data: List of data dicts
            x_key: Key for x-axis (categories)
            y_keys: Keys for y-axis (series)
            title: Chart title
            position: (left, top, width, height) in inches
            smooth: Use smooth lines

        Returns:
            Chart object
        """
        if position is None:
            position = (0.5, 1.5, 12, 5)

        left, top, width, height = [Inches(x) for x in position]

        # Chart type
        chart_type = XL_CHART_TYPE.LINE if smooth else XL_CHART_TYPE.LINE

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = [item[x_key] for item in data]

        # Add series
        for _i, y_key in enumerate(y_keys):
            values = [item[y_key] for item in data]
            chart_data.add_series(y_key, values)

        # Add chart
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart

        # Style
        self._style_chart(chart, title, len(y_keys))

        # Make lines smooth
        if smooth:
            for series in chart.series:
                series.smooth = True

        return chart

    def embed_pie_chart(
        self,
        slide,
        data: list[dict[str, Any]],
        label_key: str,
        value_key: str,
        title: str | None = None,
        position: tuple[float, float, float, float] = None,
        donut: bool = False,
    ):
        """
        Embed native pie chart in slide.

        Args:
            slide: PowerPoint slide object
            data: List of data dicts
            label_key: Key for labels
            value_key: Key for values
            title: Chart title
            position: (left, top, width, height) in inches
            donut: Whether to make donut chart

        Returns:
            Chart object
        """
        if position is None:
            position = (0.5, 1.5, 12, 5)

        left, top, width, height = [Inches(x) for x in position]

        # Chart type
        chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = [item[label_key] for item in data]

        # Add series (pie only has one)
        values = [item[value_key] for item in data]
        chart_data.add_series("Values", values)

        # Add chart
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart

        # Style
        self._style_chart(chart, title, len(data))

        # Show data labels with percentages
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        # Data labels
        if chart.plots[0].has_data_labels:
            data_labels = chart.plots[0].data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        return chart

    def embed_area_chart(
        self,
        slide,
        data: list[dict[str, Any]],
        x_key: str,
        y_keys: list[str],
        title: str | None = None,
        position: tuple[float, float, float, float] = None,
        stacked: bool = False,
    ):
        """
        Embed native area chart in slide.

        Args:
            slide: PowerPoint slide object
            data: List of data dicts
            x_key: Key for x-axis
            y_keys: Keys for y-axis
            title: Chart title
            position: Position tuple
            stacked: Whether to stack areas

        Returns:
            Chart object
        """
        if position is None:
            position = (0.5, 1.5, 12, 5)

        left, top, width, height = [Inches(x) for x in position]

        # Chart type
        chart_type = XL_CHART_TYPE.AREA_STACKED if stacked else XL_CHART_TYPE.AREA

        # Create data
        chart_data = CategoryChartData()
        chart_data.categories = [item[x_key] for item in data]

        for y_key in y_keys:
            values = [item[y_key] for item in data]
            chart_data.add_series(y_key, values)

        # Add chart
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart

        # Style
        self._style_chart(chart, title, len(y_keys))

        return chart

    def _style_chart(self, chart, title: str | None, num_series: int):
        """
        Apply KDS styling to chart.

        Args:
            chart: Chart object
            title: Chart title
            num_series: Number of data series
        """
        # Title
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

            # Title font
            title_para = chart.chart_title.text_frame.paragraphs[0]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.name = "Inter"
            title_para.font.color.rgb = RGBColor.from_string(
                self.theme.get_text().lstrip("#")
            )

        # Legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        # Legend font
        legend_font = chart.legend.font
        legend_font.size = Pt(10)
        legend_font.name = "Inter"
        legend_font.color.rgb = RGBColor.from_string(
            self.theme.get_text().lstrip("#")
        )

        # Apply KDS colors to series
        for i, series in enumerate(chart.series):
            # Get color from KDS palette
            color_hex = self.colors[i % len(self.colors)]
            rgb = RGBColor.from_string(color_hex.lstrip("#"))

            # Apply to series
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = rgb
            except Exception:
                pass  # Some chart types don't support this

        # Remove gridlines (KDS rule)
        try:
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.has_minor_gridlines = False
        except Exception:
            pass

        # Axis fonts
        try:
            # Value axis
            value_axis = chart.value_axis
            value_axis.tick_labels.font.size = Pt(10)
            value_axis.tick_labels.font.name = "Inter"
            value_axis.tick_labels.font.color.rgb = RGBColor.from_string(
                self.theme.get_text("secondary").lstrip("#")
            )

            # Category axis
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(10)
            category_axis.tick_labels.font.name = "Inter"
            category_axis.tick_labels.font.color.rgb = RGBColor.from_string(
                self.theme.get_text("secondary").lstrip("#")
            )
        except Exception:
            pass

        # Background
        try:
            chart.chart_area.fill.solid()
            chart.chart_area.fill.fore_color.rgb = RGBColor.from_string(
                self.theme.get_background().lstrip("#")
            )

            # Plot area background
            chart.plot_area.fill.solid()
            chart.plot_area.fill.fore_color.rgb = RGBColor.from_string(
                self.theme.get_background().lstrip("#")
            )
        except Exception:
            pass

    def embed_table(
        self,
        slide,
        data: list[dict[str, Any]],
        columns: list[str],
        position: tuple[float, float, float, float] = None,
    ):
        """
        Embed native table in slide.

        Args:
            slide: PowerPoint slide object
            data: List of data dicts
            columns: Column keys to include
            position: Position tuple

        Returns:
            Table object
        """
        if position is None:
            position = (0.5, 1.5, 12, 5)

        left, top, width, height = [Inches(x) for x in position]

        # Create table
        rows = len(data) + 1  # +1 for header
        cols = len(columns)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Header row
        for col_idx, col_key in enumerate(columns):
            cell = table.cell(0, col_idx)
            cell.text = col_key.replace("_", " ").title()

            # Style header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("7823DC")  # KDS purple

            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(12)
            para.font.name = "Inter"
            para.font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Data rows
        for row_idx, item in enumerate(data, start=1):
            for col_idx, col_key in enumerate(columns):
                cell = table.cell(row_idx, col_idx)
                value = item.get(col_key, "")
                cell.text = str(value)

                # Style data cell
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.font.name = "Inter"
                para.font.color.rgb = RGBColor.from_string(
                    self.theme.get_text().lstrip("#")
                )

                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        self.theme.get_background("secondary").lstrip("#")
                    )

        return table


def embed_chart_in_slide(
    prs: Presentation,
    slide_idx: int,
    chart_type: str,
    data: list[dict[str, Any]],
    **kwargs,
):
    """
    Helper function to embed chart in specific slide.

    Args:
        prs: Presentation object
        slide_idx: Slide index
        chart_type: Type of chart ('bar', 'line', 'pie', 'area')
        data: Chart data
        **kwargs: Chart-specific parameters

    Returns:
        Chart object
    """
    slide = prs.slides[slide_idx]
    embedder = PowerPointChartEmbedder()

    if chart_type == "bar":
        return embedder.embed_bar_chart(slide, data, **kwargs)
    elif chart_type == "line":
        return embedder.embed_line_chart(slide, data, **kwargs)
    elif chart_type == "pie":
        return embedder.embed_pie_chart(slide, data, **kwargs)
    elif chart_type == "area":
        return embedder.embed_area_chart(slide, data, **kwargs)
    else:
        raise ValueError(f"Unknown chart type: {chart_type}")
